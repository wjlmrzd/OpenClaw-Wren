#!/usr/bin/env python
"""
ppt-maker: 精美PPT制作工具
支持科技风设计、图文混排、HTML内容嵌入
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from PIL import Image
import os
import sys
import argparse
from enum import Enum


class Theme(Enum):
    """主题风格"""
    TECH = "tech"        # 科技风
    MODERN = "modern"    # 现代简约
    CORPORATE = "corporate"  # 企业风


class PresentationBuilder:
    """PPT构建器"""
    
    def __init__(self, theme=Theme.TECH):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.theme = theme
        self._apply_theme()
    
    def _apply_theme(self):
        """应用主题颜色"""
        if self.theme == Theme.TECH:
            self.colors = {
                'primary': RGBColor(0, 122, 204),      # 科技蓝
                'secondary': RGBColor(0, 61, 112),     # 深蓝
                'accent': RGBColor(0, 255, 255),       # 青色
                'dark': RGBColor(20, 30, 50),          # 深色背景
                'light': RGBColor(240, 250, 255),      # 浅色背景
                'text': RGBColor(255, 255, 255),       # 白色文字
                'text_dark': RGBColor(50, 50, 50),     # 深色文字
                'gradient_start': RGBColor(0, 80, 160),
                'gradient_end': RGBColor(0, 30, 60),
            }
        elif self.theme == Theme.MODERN:
            self.colors = {
                'primary': RGBColor(50, 50, 50),
                'secondary': RGBColor(100, 100, 100),
                'accent': RGBColor(255, 100, 0),
                'dark': RGBColor(30, 30, 30),
                'light': RGBColor(250, 250, 250),
                'text': RGBColor(255, 255, 255),
                'text_dark': RGBColor(50, 50, 50),
                'gradient_start': RGBColor(60, 60, 60),
                'gradient_end': RGBColor(30, 30, 30),
            }
        else:  # CORPORATE
            self.colors = {
                'primary': RGBColor(0, 102, 204),
                'secondary': RGBColor(0, 51, 102),
                'accent': RGBColor(255, 153, 0),
                'dark': RGBColor(240, 240, 240),
                'light': RGBColor(255, 255, 255),
                'text': RGBColor(255, 255, 255),
                'text_dark': RGBColor(50, 50, 50),
                'gradient_start': RGBColor(0, 102, 204),
                'gradient_end': RGBColor(0, 51, 102),
            }
    
    def _add_gradient_background(self, slide):
        """添加渐变背景"""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['dark']
        shape.line.fill.background()
    
    def _add_top_bar(self, slide, height=Inches(1.2)):
        """添加顶部色带"""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['primary']
        shape.line.fill.background()
    
    def _add_title_text(self, slide, text, x, y, width, height, font_size=36, bold=True, color=None):
        """添加标题文本"""
        box = slide.shapes.add_textbox(x, y, width, height)
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or self.colors['text']
        p.alignment = PP_ALIGN.LEFT
        return box
    
    def add_title_slide(self, title, subtitle="", bg_image=None):
        """添加标题页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景
        self._add_gradient_background(slide)
        
        # 添加装饰线条
        for i in range(5):
            y = Inches(1.5 + i * 0.15)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-2), y, Inches(3), Inches(0.05))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors['accent']
            shape.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = self.colors['text']
        p.alignment = PP_ALIGN.LEFT
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = self.colors['accent']
            p.alignment = PP_ALIGN.LEFT
        
        # 底部信息
        if self.theme == Theme.TECH:
            self._add_tech_decoration(slide)
        
        return slide
    
    def _add_tech_decoration(self, slide):
        """添加科技感装饰"""
        # 底部装饰线
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, self.prs.slide_height - Inches(0.1), self.prs.slide_width, Inches(0.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['accent']
        shape.line.fill.background()
        
        # 右侧装饰点
        for i in range(8):
            x = self.prs.slide_width - Inches(0.5 + i * 0.3)
            y = self.prs.slide_height - Inches(0.8 - i * 0.08)
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.15), Inches(0.15))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors['accent']
            shape.line.fill.background()
    
    def add_content_slide(self, title, bullets, icon="", description=""):
        """添加内容页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景
        self._add_gradient_background(slide)
        
        # 顶部色带
        self._add_top_bar(slide, Inches(1.0))
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icon} {title}" if icon else title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.colors['text']
        
        # 描述
        if description:
            desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.4))
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = description
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(180, 200, 220)
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(22)
            p.font.color.rgb = self.colors['text']
            p.space_after = Pt(18)
            # 首行缩进
            p.level = 0
        
        return slide
    
    def add_feature_grid(self, features, title="核心功能"):
        """添加特性网格页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景
        self._add_gradient_background(slide)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Star] {title}"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.colors['text']
        
        # 特性卡片
        cols = 3
        card_width = Inches(4)
        card_height = Inches(2.2)
        start_x = Inches(0.5)
        start_y = Inches(1.3)
        gap_x = Inches(0.25)
        gap_y = Inches(0.25)
        
        for i, (feat_title, feat_desc) in enumerate(features):
            row = i // cols
            col = i % cols
            x = start_x + col * (card_width + gap_x)
            y = start_y + row * (card_height + gap_y)
            
            # 卡片背景（科技风玻璃效果）
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 50, 80)
            shape.line.color.rgb = self.colors['accent']
            shape.line.width = Pt(1)
            
            # 特性图标/编号
            icon_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5))
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.colors['accent']
            
            # 特性标题
            feat_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), card_width - Inches(0.4), Inches(0.5))
            tf = feat_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = feat_title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.colors['text']
            
            # 特性描述
            desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.1), card_width - Inches(0.4), Inches(1))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = feat_desc
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(180, 200, 220)
        
        return slide
    
    def add_comparison_slide(self, title, left_title, left_items, right_title, right_items):
        """添加对比页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景
        self._add_gradient_background(slide)
        
        # 标题
        self._add_top_bar(slide, Inches(0.9))
        self._add_title_text(slide, f"[VS] {title}", Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), font_size=32)
        
        # 左侧
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4))
        tf = left_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[X] {left_title}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 100, 100)
        
        left_content = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5.4), Inches(5))
        tf = left_content.text_frame
        for i, item in enumerate(left_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(200, 180, 180)
            p.space_after = Pt(12)
        
        # 右侧
        right_box = slide.shapes.add_textbox(Inches(7), Inches(1.2), Inches(5.8), Inches(0.4))
        tf = right_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[OK] {right_title}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.colors['accent']
        
        right_content = slide.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.4), Inches(5))
        tf = right_content.text_frame
        for i, item in enumerate(right_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['text']
            p.space_after = Pt(12)
        
        # 中间分割线
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.2), Inches(0.02), Inches(6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors['accent']
        shape.line.fill.background()
        
        return slide
    
    def add_image_slide(self, title, image_path, description=""):
        """添加图片页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景
        self._add_gradient_background(slide)
        
        # 标题
        self._add_top_bar(slide)
        self._add_title_text(slide, title, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6), font_size=36)
        
        # 图片
        if os.path.exists(image_path):
            # 计算图片尺寸，保持比例
            img = Image.open(image_path)
            img_width, img_height = img.size
            max_width = Inches(12)
            max_height = Inches(5.5)
            
            ratio = min(max_width.inches / img_width, max_height.inches / img_height)
            final_width = img_width * ratio
            final_height = img_height * ratio
            
            x = (self.prs.slide_width - final_width) / 2
            y = Inches(1.5)
            
            slide.shapes.add_picture(image_path, x, y, width=final_width)
        
        # 描述
        if description:
            desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.333), Inches(0.5))
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = description
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(180, 200, 220)
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_html_snapshot(self, title, html_content, description=""):
        """
        添加HTML内容（通过内嵌文本模拟）
        注意：完整的HTML嵌入需要额外处理，这里提供文本摘要展示
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景
        self._add_gradient_background(slide)
        
        # 标题
        self._add_top_bar(slide)
        self._add_title_text(slide, f"[Web] {title}", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6), font_size=32)
        
        # 代码/内容框
        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5))
        tf = code_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        # 截取HTML内容片段展示
        preview = html_content[:500] + "..." if len(html_content) > 500 else html_content
        p.text = preview
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(0, 255, 180)
        
        if description:
            desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = description
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(180, 200, 220)
        
        return slide
    
    def add_closing_slide(self, title, subtitle=""):
        """添加结束页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景
        self._add_gradient_background(slide)
        
        # 装饰
        for i in range(10):
            x = Inches(i * 1.5)
            y = self.prs.slide_height - Inches(0.5)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.8), Inches(0.05))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors['accent']
            shape.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(0, Inches(3), self.prs.slide_width, Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.colors['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(0, Inches(4.2), self.prs.slide_width, Inches(0.8))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = self.colors['accent']
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def save(self, filename):
        """保存PPT"""
        self.prs.save(filename)
        print(f"[OK] PPT已保存: {filename}")


def create_openclow_presentation():
    """创建OpenClaw介绍PPT（科技风）"""
    prs = PresentationBuilder(theme=Theme.TECH)
    
    # 1. 封面
    prs.add_title_slide(
        "OpenClaw",
        "您的跨平台AI个人助理 | 让AI成为你的第二大脑"
    )
    
    # 2. 什么是OpenClaw
    prs.add_content_slide(
        "什么是 OpenClaw?",
        [
            "[Run] 开源免费的自托管 AI 网关",
            "[Mobile] 连接 WhatsApp、Telegram、Discord、iMessage 等多平台",
            "🔒 数据完全掌控在自己手中，不依赖第三方服务",
            "[AI] 内置编码 Agent，支持工具调用、会话管理、多 Agent 路由",
            "[Web] 在任意设备上通过消息应用与 AI 助手对话"
        ],
        icon="[AI]",
        description="一个Gateway同时连接多个通道，数据本地处理，安全可控"
    )
    
    # 3. 核心功能
    prs.add_feature_grid([
        ("多通道网关", "一个 Gateway 同时连接 WhatsApp、Telegram、Discord 等多个平台"),
        ("插件扩展", "支持 Mattermost、Feishu 等更多平台插件"),
        ("多 Agent 路由", "隔离的会话空间，支持多 Agent 协作和负载均衡"),
        ("媒体支持", "发送接收图片、音频、文件，完整多媒体体验"),
        ("Web 控制面板", "浏览器Dashboard 管理配置、会话和节点"),
        ("移动节点", "配对 iOS/Android 设备，支持 Canvas 交互"),
        ("定时任务", "支持 Cron 定时执行任务和提醒"),
        ("智能记忆", "持久化会话上下文，理解对话历史"),
        ("安全控制", "Token 认证、IP 白名单、操作审计")
    ], title="核心功能")
    
    # 4. 效率提升对比
    prs.add_comparison_slide(
        "效率提升对比",
        "传统方式", [
            "手动打开各个App查看消息",
            "在不同平台间切换操作",
            "重复复制粘贴信息",
            "无法自动化处理任务",
            "消息分散难以管理"
        ],
        "使用 OpenClaw", [
            "统一消息入口，7×24自动响应",
            "一个界面管理所有渠道",
            "自动提取和整理信息",
            "定时任务+自动化工作流",
            "会话持久化，随时继续"
        ]
    )
    
    # 5. 实际应用场景
    prs.add_content_slide(
        "实际应用场景",
        [
            "[Chat] 微信/QQ 自动回复：7×24小时响应，不错过任何消息",
            "[Cal] 日程管理：自动检查日历，重要事件提前提醒",
            "[Search] 信息聚合：跨平台搜索，统一整理工作资料",
            "[Doc] 内容发布：一键发布到公众号、小红书等平台",
            "[PC] 远程控制：通过消息执行终端命令，管理服务器",
            "[Bell] 智能提醒：定时任务触发通知，支持多渠道"
        ],
        icon="[Run]",
        description="从简单自动回复到复杂工作流，释放双手"
    )
    
    # 6. 技术架构
    prs.add_content_slide(
        "技术架构",
        [
            "[Arch] Gateway 模式：单一进程处理所有通道连接和路由",
            "[Plug] MCP 协议：标准化工具调用，支持 Python/Node 扩展",
            "[Folder] Session 管理：每个对话独立隔离，支持上下文持久化",
            "[Shield] 安全控制：Token 认证、IP 白名单、操作审计",
            "[PC] 跨平台支持：Windows、macOS、Linux、VPS 都能运行",
            "[Sync] 插件系统：灵活的扩展机制，轻松添加新功能"
        ],
        icon="[Arch]",
        description="模块化设计，灵活可扩展"
    )
    
    # 7. 快速开始
    prs.add_content_slide(
        "快速开始",
        [
            "1. 安装：npm install -g openclaw@latest",
            "2. 初始化：openclaw onboard --install-daemon",
            "3. 登录通道：openclaw channels login",
            "4. 启动网关：openclaw gateway --port 18789",
            "5. 打开控制台：http://127.0.0.1:18789"
        ],
        icon="[List]",
        description="5分钟快速上手"
    )
    
    # 8. 结束页
    prs.add_closing_slide(
        "让 AI 成为你的第二大脑",
        "用 OpenClaw 打造你的专属 AI 助手"
    )
    
    # 保存
    output_path = r"C:\Users\90781\Desktop\OpenClaw科技风.pptx"
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPT制作工具")
    parser.add_argument("--title", "-t", help="标题")
    parser.add_argument("--content", "-c", help="内容（用|分隔）")
    parser.add_argument("--output", "-o", help="输出文件")
    parser.add_argument("--theme", default="tech", choices=["tech", "modern", "corporate"])
    parser.add_argument("--openclaw", action="store_true", help="生成OpenClaw介绍PPT")
    
    args = parser.parse_args()
    
    if args.openclaw:
        output = create_openclow_presentation()
        print(f"\n[Done] OpenClaw 科技风 PPT 已生成: {output}")
    elif args.title and args.content:
        theme = Theme.TECH if args.theme == "tech" else Theme.MODERN if args.theme == "modern" else Theme.CORPORATE
        prs = PresentationBuilder(theme=theme)
        bullets = args.content.split("|")
        prs.add_content_slide(args.title, bullets)
        output = args.output or "output.pptx"
        prs.save(output)
        print(f"\n[Done] PPT已生成: {output}")
    else:
        print("使用 --openclaw 生成OpenClaw介绍，或使用 --title --content 指定内容")
        print("示例: python ppt_maker.py -t '标题' -c '要点1|要点2|要点3' -o demo.pptx")
